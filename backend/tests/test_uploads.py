"""Upload ingestion tests — extraction + LLM 判讀 across file types."""
from __future__ import annotations

import io

import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from app.models import Document, DocumentStatus, GuardrailEvent, GuardrailType
from app.services.extraction import ExtractionError, extract_text
from app.services.ingestion import ingest_file

PROPOSAL = "目標:提升風控。範圍:全行。時程:Q3。風險:資料缺漏。資源:兩名工程師。里程碑:三階段。"


def _docx_bytes(text: str) -> bytes:
    doc = DocxDocument()
    for line in text.split("。"):
        if line:
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(text: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 600, 400)
    box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pdf_bytes(text: str) -> bytes:
    """Hand-assemble a minimal one-page PDF with a Helvetica text run.

    Avoids a PDF-writing dependency just for test fixtures — extraction.py
    only needs to *read* PDFs (via pypdf), so there's no writer available.
    """
    escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
    content = f"BT /F1 24 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1")
    stream_obj = f"<< /Length {len(content)} >>\nstream\n".encode("latin-1") + content + b"\nendstream"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        stream_obj,
    ]
    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = [0]
    for i, body in enumerate(objects, start=1):
        offsets.append(out.tell())
        out.write(f"{i} 0 obj\n".encode("latin-1"))
        out.write(body)
        out.write(b"\nendobj\n")
    xref_offset = out.tell()
    n = len(objects) + 1
    out.write(f"xref\n0 {n}\n".encode("latin-1"))
    out.write(b"0000000000 65535 f \n")
    for off in offsets[1:]:
        out.write(f"{off:010d} 00000 n \n".encode("latin-1"))
    out.write(b"trailer\n")
    out.write(f"<< /Size {n} /Root 1 0 R >>\n".encode("latin-1"))
    out.write(b"startxref\n")
    out.write(f"{xref_offset}\n".encode("latin-1"))
    out.write(b"%%EOF")
    return out.getvalue()


def test_extract_docx():
    kind, text = extract_text("plan.docx", _docx_bytes(PROPOSAL))
    assert kind == "docx"
    assert "目標" in text and "里程碑" in text


def test_extract_pptx():
    kind, text = extract_text("deck.pptx", _pptx_bytes(PROPOSAL))
    assert kind == "pptx"
    assert "Slide 1" in text and "目標" in text


def test_extract_pdf():
    # The hand-rolled fixture only embeds the base Helvetica font (no CJK
    # glyphs), so it uses ASCII; real-world PDFs with CJK extract fine via
    # pypdf as long as the PDF itself embeds a CJK-capable font.
    kind, text = extract_text("plan.pdf", _pdf_bytes("Goal Scope Timeline Risk Milestone"))
    assert kind == "pdf"
    assert "Milestone" in text


def test_unsupported_type_raises():
    with pytest.raises(ExtractionError):
        extract_text("image.png", b"\x89PNG")


def test_ingest_complete_proposal_is_green(db):
    doc = ingest_file(db, filename="plan.docx", data=_docx_bytes(PROPOSAL))
    assert doc.status == DocumentStatus.DONE
    assert doc.char_count > 0
    assert doc.llm_verdict == "綠燈"


def test_ingest_failure_records_capability_guardrail(db):
    doc = ingest_file(db, filename="broken.docx", data=b"not a real docx")
    assert doc.status == DocumentStatus.FAILED
    assert doc.error
    events = db.query(GuardrailEvent).filter_by(guardrail_type=GuardrailType.CAPABILITY).all()
    assert len(events) == 1


def test_ingest_multiple_kinds(db):
    # The hand-rolled PDF fixture only embeds Helvetica (no CJK glyphs), so
    # it gets ASCII text — see test_extract_pdf for why.
    files = {
        "a.pdf": _pdf_bytes("Goal Scope Timeline Risk Milestone"),
        "b.pptx": _pptx_bytes(PROPOSAL),
        "c.docx": _docx_bytes(PROPOSAL),
    }
    for name, data in files.items():
        ingest_file(db, filename=name, data=data)
    assert db.query(Document).count() == 3
