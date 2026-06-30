"""Document text extraction → Markdown.

Accepts DOCX / PPTX / TXT / MD and turns each into clean **Markdown**, which
is the canonical representation everything downstream reads: the AI 評審中心
(LLM 判讀), the flowchart generator and the structured-field extractor all
consume this Markdown rather than a flat text blob. Producing Markdown up front
("資料分類做成 md 再讀取") keeps headings, tables and slide structure intact, so
the LLM sees a well-formed document instead of a wall of text.

PDF is intentionally not supported: every PDF-parsing library tried so far
(PyMuPDF, pypdf) got blocked by the company's package-approval policy, which
applies to any new dependency regardless of license. Re-add PDF support once
a library has been approved, or implement a pure-stdlib parser if that
approval never comes.

Per the Capability guardrail (§八), a parse failure must surface as an
explicit "無法取得資料" rather than fabricated content — so callers treat an
``ExtractionError`` as a hard stop, never as empty-but-ok text.
"""
from __future__ import annotations

import io

from docx import Document as DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation


class ExtractionError(Exception):
    """Raised when a file cannot be parsed (→ Capability guardrail)."""


# extension → canonical kind
SUPPORTED = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "txt",
}


def detect_kind(filename: str) -> str | None:
    lower = filename.lower()
    for ext, kind in SUPPORTED.items():
        if lower.endswith(ext):
            return kind
    return None


def _heading_level(style_name: str | None) -> int | None:
    """Map a docx paragraph style to a Markdown heading level (1..6)."""
    if not style_name:
        return None
    name = style_name.strip()
    if name.lower() in ("title",):
        return 1
    if name.startswith("Heading"):
        try:
            level = int(name.split()[-1])
        except (ValueError, IndexError):
            return 2
        return min(max(level, 1), 6)
    return None


def _is_bullet(style_name: str | None) -> bool:
    if not style_name:
        return False
    lowered = style_name.lower()
    return "list" in lowered or "bullet" in lowered


def _table_to_markdown(table: Table) -> list[str]:
    """Render a docx table as a GitHub-flavoured Markdown table."""
    rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
    rows = [r for r in rows if any(c for c in r)]
    if not rows:
        return []
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return out


def _extract_docx(data: bytes) -> str:
    doc = DocxDocument(io.BytesIO(data))
    parts: list[str] = []
    body = doc.element.body
    # Walk the body in document order so headings, paragraphs and tables keep
    # their original sequence (python-docx exposes paragraphs and tables in
    # separate collections, which loses interleaving).
    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            level = _heading_level(para.style.name if para.style else None)
            if level:
                parts.append(f"{'#' * level} {text}")
            elif _is_bullet(para.style.name if para.style else None):
                parts.append(f"- {text}")
            else:
                parts.append(text)
        elif child.tag.endswith("}tbl"):
            table = Table(child, doc)
            md = _table_to_markdown(table)
            if md:
                parts.append("")
                parts.extend(md)
                parts.append("")
    return "\n".join(parts).strip()


def _extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    parts.append(f"- {text}")
        parts.append("")
    return "\n".join(parts).strip()


def extract_text(filename: str, data: bytes) -> tuple[str, str]:
    """Return (kind, markdown). Raises ExtractionError on failure.

    The second element is Markdown (TXT/MD are passed through as-is, since they
    are already plain/markdown text).
    """
    kind = detect_kind(filename)
    if kind is None:
        raise ExtractionError(f"不支援的檔案格式:{filename}")
    if not data:
        raise ExtractionError("檔案內容為空,無法解析。")
    try:
        if kind == "docx":
            text = _extract_docx(data)
        elif kind == "pptx":
            text = _extract_pptx(data)
        else:  # txt / md
            text = data.decode("utf-8", errors="replace")
    except ExtractionError:
        raise
    except Exception as exc:  # noqa: BLE001 — any parser error → Capability halt
        raise ExtractionError(f"解析失敗:{exc}") from exc

    if not text.strip():
        raise ExtractionError("檔案中未擷取到任何文字內容。")
    return kind, text
