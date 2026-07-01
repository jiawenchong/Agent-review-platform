"""產生「AI Agent 驗證報告」PPTX — 深色科技風格,10 頁固定結構。

用法:
    python generate_ppt.py [input_data.json] [output.pptx]

    input_data.json  — 選填,內容欄位對照 reference/input_data.example.json;
                        省略則產出全部留白佔位符(方括號提示文字)的空白模板。
    output.pptx       — 選填,預設 AI_Agent_Validation_Report.pptx(存在目前工作目錄)。

頁面結構、配色與字體規則的權威定義在 ../reference/template_structure.json
(single source of truth) — 本腳本在啟動時讀取它,不要把版面規則寫死兩份。
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

_SKILL_ROOT = Path(__file__).resolve().parents[1]
_TEMPLATE_STRUCTURE_PATH = _SKILL_ROOT / "reference" / "template_structure.json"

with open(_TEMPLATE_STRUCTURE_PATH, "r", encoding="utf-8") as f:
    TEMPLATE_CONFIG = json.load(f)

_META = TEMPLATE_CONFIG["template_metadata"]
_ACCENT = _META["accent_colors"]


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


COLORS = {
    "bg": _rgb(_META["background_color"]),
    "title": _rgb(_ACCENT["text_main"]),
    "body": _rgb(_ACCENT["text_sub"]),
    "accent": _rgb(_ACCENT["primary"]),
    "secondary": _rgb(_ACCENT["secondary"]),
    "warning": _rgb(_ACCENT["warning"]),
}
# Table cells default to a white fill in python-pptx regardless of the slide
# background, which would make the light-gray/green theme text unreadable —
# every table cell gets an explicit dark navy fill to match the deck.
CELL_BG = RGBColor(0x14, 0x1B, 0x33)
HEADER_CELL_BG = RGBColor(0x1A, 0x2E, 0x4A)

TITLE_FONT = _META["fonts"]["title"]
BODY_FONT = _META["fonts"]["body"]


def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs: Presentation):
    """Blank slide with the dark theme background applied.

    Note: python-pptx has no presentation-level `slide_background` — the
    background must be set per slide via `slide.background`.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]
    return slide


def add_textbox(slide, left, top, width, height, text, font_size=14, color=None,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    color = color or COLORS["body"]
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name or BODY_FONT
        p.alignment = alignment
    return box


def add_table(slide, left, top, width, height, data, col_widths=None):
    """data: list of rows (list of cell text); first row is treated as header."""
    rows, cols = len(data), len(data[0])
    frame = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = frame.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_CELL_BG if i == 0 else CELL_BG
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.name = BODY_FONT
                para.font.bold = i == 0
                para.font.color.rgb = COLORS["accent"] if i == 0 else COLORS["body"]
    return frame


# ── per-page generators (order matches reference/template_structure.json) ──


def generate_slide_cover(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(2), Inches(1), "[公司 Logo]", font_size=10)
    add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(1.5), data.get("title", "[報告主題]"),
                font_size=44, color=COLORS["title"], bold=True, alignment=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    add_textbox(slide, Inches(2), Inches(4.0), Inches(9), Inches(1), data.get("subtitle", "[專案名稱]"),
                font_size=24, color=COLORS["accent"], alignment=PP_ALIGN.CENTER)
    info = data.get("info", {})
    info_text = (
        f"報告人: {info.get('reporter', '[報告人姓名]')}\n"
        f"單位: {info.get('department', '[單位名稱]')}\n"
        f"日期: {info.get('date', '[報告日期]')}"
    )
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(4), Inches(1.5), info_text, font_size=14)
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.5), "ASE Confidential / Security-C",
                font_size=10, alignment=PP_ALIGN.CENTER)


def generate_slide_goal_definition(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "Agent 目標定義",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)

    left_table = [
        ["欄位", "內容"],
        ["專案說明", data.get("project_desc", "[專案說明]")],
        ["Agent 角色", data.get("agent_role", "[Agent 角色]")],
        ["核心使命", data.get("mission", "[核心使命]")],
        ["觸發時機", data.get("trigger", "[觸發時機]")],
    ]
    add_table(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(3), left_table, col_widths=[2, 3])

    add_textbox(slide, Inches(6), Inches(1.5), Inches(6), Inches(0.5), "Agent 指標",
                font_size=20, color=COLORS["accent"], bold=True)
    add_textbox(
        slide, Inches(6), Inches(2.2), Inches(6), Inches(2),
        f"成功門檻: {data.get('success_threshold', '[成功門檻]')}\n\n"
        f"專案指標: {data.get('project_metrics', '[專案指標]')}",
        font_size=14,
    )
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1),
                "流程: Trigger → Perception → Reasoning → Action → Feedback",
                font_size=14, color=COLORS["secondary"], alignment=PP_ALIGN.CENTER)


def generate_slide_data_sources(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "資料來源與工具授權",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)

    col_titles = ["資料來源 (Data)", "模型使用 (Model)", "知識庫內容 (Knowledge)"]
    col_contents = [
        data.get("data_sources", "[資料來源內容]"),
        data.get("model_usage", "[模型使用內容]"),
        data.get("knowledge_base", "[知識庫內容]"),
    ]
    for i, (title, content) in enumerate(zip(col_titles, col_contents)):
        left = Inches(0.5 + i * 4.2)
        add_textbox(slide, left, Inches(1.5), Inches(3.8), Inches(0.5), title,
                    font_size=18, color=COLORS["accent"], bold=True)
        add_textbox(slide, left, Inches(2.2), Inches(3.8), Inches(2.5), content, font_size=14)

    table_data = [
        ["Skills", "Tools", "Data", "Source"],
        [data.get("skills", "[Skills]"), data.get("tools", "[Tools]"),
         data.get("data", "[Data]"), data.get("source", "[Source]")],
    ]
    add_table(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(2), table_data, col_widths=[3, 3, 3, 3])


def generate_slide_decision_logic(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "決策思維與流程",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)

    table_data = [
        ["Tasks", "Sub-Agent", "思考邏輯", "Tools & Skills"],
        [data.get("tasks", "[Tasks]"), data.get("sub_agent", "[Sub-Agent]"),
         data.get("logic", "[思考邏輯]"), data.get("logic_tools", "[Tools & Skills]")],
    ]
    add_table(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(3), table_data, col_widths=[1.5, 1.5, 1.5, 1.5])

    add_textbox(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.5), "Decision Flow",
                font_size=20, color=COLORS["accent"], bold=True)
    add_textbox(
        slide, Inches(7), Inches(2.2), Inches(5), Inches(2),
        "[流程圖佔位符]\n1. 進度是否落後?\n2. 負責人是否有填寫 Recovery Plan?\n"
        "3. GPT 判斷合理性\n4. Agenda Priority #1~#4",
        font_size=14,
    )
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(6), Inches(0.5),
                "註: 正式派發三十分鐘內,重填或補填", font_size=12, color=COLORS["warning"])


def generate_slide_guardrails(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "產線防護與紅線設定",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)

    table_data = [
        ["No", "Guardrails Lists", "Description"],
        ["1", "Output", data.get("g1_output", "[Output 描述]")],
        ["2", "Capability", data.get("g2_capability", "[Capability 描述]")],
        ["3", "Grounding", data.get("g3_grounding", "[Grounding 描述]")],
        ["4", "Hallucination", data.get("g4_hallucination", "[Hallucination 描述]")],
        ["5", "Goal", data.get("g5_goal", "[Goal 描述]")],
    ]
    add_table(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), table_data, col_widths=[1, 2, 9])


def generate_slide_golden_tests(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "黃金測試題庫",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)

    table_data = [
        ["Scenario", "問題狀況", "專家標準決策", "Agent 內容驗證檢查點"],
        ["進度落後無 Plan", data.get("t1_problem", "[問題狀況]"), data.get("t1_expert", "[專家決策]"), data.get("t1_check", "[檢查點]")],
        ["理由不合理", data.get("t2_problem", "[問題狀況]"), data.get("t2_expert", "[專家決策]"), data.get("t2_check", "[檢查點]")],
        ["數據矛盾", data.get("t3_problem", "[問題狀況]"), data.get("t3_expert", "[專家決策]"), data.get("t3_check", "[檢查點]")],
        ["檔案損壞", data.get("t4_problem", "[問題狀況]"), data.get("t4_expert", "[專家決策]"), data.get("t4_check", "[檢查點]")],
    ]
    add_table(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5), table_data, col_widths=[2, 3, 3.5, 3.5])


def generate_slide_task_capability(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "任務完成能力",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                "成功標準: 任務完成率 100%,結果準確率 100%", font_size=14)

    add_table(slide, Inches(0.5), Inches(2.0), Inches(3), Inches(1),
              [["指標", "數值"], ["任務成功率", data.get("task_rate", "[數據]")]], col_widths=[1.5, 1.5])
    add_table(slide, Inches(0.5), Inches(3.2), Inches(3), Inches(1),
              [["指標", "數值"], ["結果準確率", data.get("accuracy_rate", "[數據]")]], col_widths=[1.5, 1.5])

    add_textbox(slide, Inches(4.5), Inches(2.0), Inches(8), Inches(0.5), "驗證截圖",
                font_size=18, color=COLORS["accent"], bold=True)
    add_textbox(
        slide, Inches(4.5), Inches(2.7), Inches(8), Inches(3),
        "[截圖 1: 工程條列]\n\n[截圖 2: 進度解析]\n\n[截圖 3: Agenda 優先順序]", font_size=12,
    )


def generate_slide_offline_validation(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "Offline 驗證",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                f"驗證天數: {data.get('days', '[天數]')} 天 | 期間: {data.get('period', '[期間]')}", font_size=14)

    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(5), Inches(0.5), "郵件清單",
                font_size=18, color=COLORS["accent"], bold=True)
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(5), Inches(3), "[郵件清單截圖佔位符]", font_size=12)

    add_textbox(slide, Inches(7), Inches(2.0), Inches(5), Inches(0.5), "每日產出正確日會 agenda",
                font_size=18, color=COLORS["accent"], bold=True)
    add_textbox(slide, Inches(7), Inches(2.7), Inches(5), Inches(3), "[截圖佔位符]\n\n成功率: 100%", font_size=12)


def generate_slide_online_validation(prs, data):
    slide = add_slide(prs)
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1), "Online 驗證",
                font_size=32, color=COLORS["title"], bold=True, font_name=TITLE_FONT)
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                f"持續運行: {data.get('run_days', '[天數]')} 天 | 每日準時完成", font_size=14)

    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(5), Inches(0.5), "郵件清單",
                font_size=18, color=COLORS["accent"], bold=True)
    add_textbox(slide, Inches(0.5), Inches(2.7), Inches(5), Inches(3), "[郵件清單截圖佔位符]", font_size=12)

    add_textbox(slide, Inches(7), Inches(2.0), Inches(5), Inches(0.5), "效能對比",
                font_size=18, color=COLORS["accent"], bold=True)
    add_textbox(
        slide, Inches(7), Inches(2.7), Inches(5), Inches(3),
        "[長條圖 1: 產出成功率 100%]\n\n[長條圖 2: 人工 30min vs Agent 3min]", font_size=12,
    )


def generate_slide_closing(prs):
    slide = add_slide(prs)
    add_textbox(slide, Inches(2), Inches(3), Inches(9), Inches(1), "Q&A",
                font_size=48, color=COLORS["title"], bold=True, alignment=PP_ALIGN.CENTER, font_name=TITLE_FONT)
    add_textbox(slide, Inches(2), Inches(4.5), Inches(9), Inches(0.5), "Thank You",
                font_size=24, color=COLORS["accent"], alignment=PP_ALIGN.CENTER)


def build(input_data: dict) -> Presentation:
    prs = create_presentation()
    generate_slide_cover(prs, input_data)
    generate_slide_goal_definition(prs, input_data)
    generate_slide_data_sources(prs, input_data)
    generate_slide_decision_logic(prs, input_data)
    generate_slide_guardrails(prs, input_data)
    generate_slide_golden_tests(prs, input_data)
    generate_slide_task_capability(prs, input_data)
    generate_slide_offline_validation(prs, input_data)
    generate_slide_online_validation(prs, input_data)
    generate_slide_closing(prs)
    return prs


def main() -> None:
    input_path = Path(sys.argv[1]) if len(sys.argv) > 1 else None
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("AI_Agent_Validation_Report.pptx")

    input_data = {}
    if input_path and input_path.exists():
        with open(input_path, "r", encoding="utf-8") as f:
            input_data = json.load(f)
    elif input_path:
        print(f"⚠ 找不到 {input_path},將產出全部留白佔位符的空白模板。")

    prs = build(input_data)
    prs.save(output_path)
    print(f"已產出:{output_path}")


if __name__ == "__main__":
    main()
