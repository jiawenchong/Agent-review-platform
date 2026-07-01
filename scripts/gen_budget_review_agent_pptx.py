"""產生「預算AI審核 Agent 開發規劃書」正式版 PPTX。

用法：python scripts/gen_budget_review_agent_pptx.py
輸出：/tmp/claude-0/-home-user-Agent-review-platform/117f2e83-6710-50e4-b25d-bed2206b9b67/scratchpad/預算AI審核Agent規劃書.pptx
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path("/tmp/claude-0/-home-user-Agent-review-platform/117f2e83-6710-50e4-b25d-bed2206b9b67/scratchpad/預算AI審核Agent規劃書.pptx")

# ─────────────────────────────────────────────────────────────────────────────
# 投影片內容（標題, [段落…]）
# ─────────────────────────────────────────────────────────────────────────────
SLIDES: list[tuple[str, list[str]]] = [
    (
        "預算AI審核 Agent 開發規劃書（GRAB v1.0）",
        [
            "治理就緒藍圖 — Governance-Ready Agent Blueprint",
            "",
            "Agent 名稱：預算AI審核 Agent",
            "提案人 / 部門：廠務部門",
            "時程：2025/05/02 ~ 2025/08/21",
            "版本：v1.0",
            "",
            "評審狀態：（由系統填寫：綠燈 / 紅燈 / 待補件）",
        ],
    ),
    (
        "一、目標（Goal）★必填",
        [
            "一句話定位：",
            "  對廠務系統所有預算案件進行 AI 初步預審，提供廠務專家結構化的審核參考建議，",
            "  縮短審核 cycle time，取代原本以 Excel 為主的人工作業流程。",
            "",
            "成功指標（可量化）：",
            "  · 治理指標：所有預算案件 100% 經 AI 初步篩選，不漏判不誤判",
            "  · 效率指標：審核 cycle time 較人工降低 >50%（目前人工平均 4 小時/件）",
            "  · 品質指標：AI 預審結論與廠務專家最終審核一致率 ≥ 85%",
            "",
            "其他效益：",
            "  · 取代 Excel 作業，建立可視化預算審核平台",
            "  · 自動派發案件給對應廠務專家，減少人工轉介",
            "  · AI 審核 API 串接，提供結構化審核建議與歷史案例引用",
        ],
    ),
    (
        "二、範圍（Scope）★必填",
        [
            "納入範圍：",
            "  · 所有廠務系統的預算案件申請",
            "    （含設備採購、維修費用、工程費用、耗材採購等類型）",
            "  · 案件初步篩選：自動判斷是否符合基本申請規範",
            "  · AI 預審建議：提供批准建議、退件原因、類似歷史案例引用",
            "  · 自動派發：依廠務系統歸屬自動指派對應廠務專家",
            "  · 可視化平台：案件狀態追蹤、審核進度監控、使用情況儀表板",
            "",
            "排除範圍：",
            "  · 非廠務系統的預算案件（其他部門系統）",
            "  · 最終審批決策（仍由廠務專家或主管負責，AI 僅提供建議）",
            "  · 財務系統的過帳與請款流程",
        ],
    ),
    (
        "三、時程與里程碑（Timeline / Milestone）★必填",
        [
            "總時程：2025/05/02 ~ 2025/08/21（約 16 週）",
            "",
            "M1  需求設計    2025/05/02 ~ 2025/05/31",
            "    交付物：系統架構設計文件、資料流程圖、UI 原型",
            "    內容：現有 Excel 流程梳理（AS IS）、廠務系統介面盤點、",
            "           RAG 知識庫資料收集規劃",
            "",
            "M2  平台開發    2025/06/01 ~ 2025/07/31",
            "    交付物：可視化審核平台（前後端）、AI 審核 API 串接完成",
            "    內容：自動派發功能、ProphetAI 預審 API 整合、",
            "           歷史案件 RAG 知識庫建置、權限卡控模組",
            "",
            "M3  驗證上線    2025/08/01 ~ 2025/08/20",
            "    交付物：UAT 測試報告、上線部署文件",
            "    內容：系統整合測試、廠務專家 UAT、效能調校",
            "",
            "全系統正式上線：2025/08/21",
        ],
    ),
    (
        "四、風險與緩解（Risk）★必填",
        [
            "風險一：廠務系統存取權限問題",
            "  影響：Agent 無法讀取預算案件資料，導致功能失效",
            "  可能性：中",
            "  緩解措施：",
            "    · 平台設計已加入角色權限卡控（RBAC），只有授權廠務人員可存取",
            "    · 使用情況監控模組：異常存取即時告警",
            "    · 不依賴 IT 帳號同步，採用獨立帳號管理機制",
            "",
            "風險二：AI 預審準確率初期不足",
            "  影響：誤判率偏高，廠務專家信任度低",
            "  緩解措施：",
            "    · 初期以「建議參考」模式運作，不取代人工決策",
            "    · 每週收集專家複審結果，自動更新 RAG 知識庫",
            "    · 目標一致率 ≥85% 後再逐步提升自動化比例",
            "",
            "風險三：廠務人員使用意願低",
            "  緩解措施：提供操作說明文件；UI 盡量簡化；初期由提案人陪同使用",
        ],
    ),
    (
        "五、資源需求（Resource）★必填",
        [
            "人力：",
            "  · 開發者 × 1（全程自主開發，不需 IT 部門支援）",
            "  · 廠務 SME（Subject Matter Expert）× 1~2（提供審核邏輯與歷史案例）",
            "",
            "運算：",
            "  · ProphetAI API（公司線上 LLM，無需額外 GPU 或雲端資源）",
            "  · 本機部署（Windows 伺服器或筆電均可，僅需 Python + Node.js）",
            "",
            "資料來源：",
            "  · 廠務系統預算申請歷史資料（過去核准/退件案例）",
            "  · 廠務工程系統歸屬對照表",
            "  · RAG 知識庫（由歷史案例自動建立，持續更新）",
            "",
            "外部依賴：",
            "  · 公司 ProphetAI API 金鑰（已申請）",
            "  · 預算案件通知管道（Email / PBMS 系統 FAC 關卡）",
        ],
    ),
    (
        "六、系統架構與模組（Architecture）★必填",
        [
            "採用 Closed Loop 四階段架構：",
            "",
            "Perception  感知 / 掃描來源",
            "  · 接收來自 PBMS 系統 FAC 關卡的預算申請（Email 通知）",
            "  · 自動下載並解析申請文件（PDF / PPT / Word / Excel）→ 結構化 Markdown",
            "  · 廠務工程系統歸屬判斷（關鍵字比對 + RAG，信心分數 ≥ 0.7）",
            "",
            "Reasoning  判斷邏輯（LLM）",
            "  · ProphetAI API 進行初步預審：",
            "    - 申請是否符合廠務系統基本規範",
            "    - 比對歷史批准/退件案例（RAG 知識庫）",
            "    - 生成結構化審核建議（批准 / 退件 / 補件）及原因",
            "",
            "Action  依判斷分級派發",
            "  · 自動派發給對應廠務系統的負責專家",
            "  · 提供 AI 審核建議報告供專家參考",
            "  · 取代 Excel 作業：可視化平台顯示案件狀態、進度",
            "",
            "Feedback  驗證迴圈",
            "  · 廠務專家複審結果自動回饋 → 更新 RAG 知識庫",
            "  · 專家與 AI 一致率統計，持續優化預審準確率",
            "  · 使用情況監控 + 異常告警",
        ],
    ),
    (
        "七、流程定義（Flow Definition）— 驅動自動流程圖",
        [
            "AS IS 現況流程（人工作業）：",
            "  · PBMS 系統收到預算申請 Email",
            "  · 廠務人員人工下載案件（Excel 記錄）",
            "  · 依過往經驗判斷歸屬廠務系統",
            "  · 人工派發給對應廠務專家",
            "  · 專家判斷申請是否有基本需求 → 不合格直接退件",
            "  · 合格：尋找歷史資料、寫批准原因 → 完成",
            "  · 不合格：寫退件原因 → 退件",
            "  · 將結果手動複製回 Excel，派發週報",
            "  → 每件平均耗時 4 小時，Y25 共 356 件 = 1,585 小時/年",
            "",
            "TO BE 目標流程（AI 代理）：",
            "S1 [start] 預算申請到FAC關卡",
            "S2 [process] 擷取文件內容 python -> S3",
            "S3 [process] 廠務系統判斷RAG python -> S4",
            "S4 [decision] 信心分數>=0.7? -> 是:S5 | 否:S8",
            "S5 [process] 分配廠務系統 GPT -> S6",
            "S6 [decision] 初步退件判斷 GPT -> 符合:S8 | 不符合:S7",
            "S7 [process] 生成退件原因 GPT -> SE",
            "S8 [process] 擷取歷史案件特徵 python -> S9",
            "S9 [decision] 內容審核 GPT -> 批准:S10 | 退件:S11",
            "S10 [process] 生成批准原因及案例引用 GPT -> SE",
            "S11 [process] 生成退件說明 GPT -> S12",
            "S12 [process] 派發廠務專家複審 RPA -> SE",
            "SE [end] 完成審核並更新平台",
        ],
    ),
    (
        "附件：評審燈號說明 & 預期效益",
        [
            "評審燈號說明：",
            "  綠燈：規劃書涵蓋 ≥4 項必要章節，內容具體可行 → 自動建立專案並進入進度監控",
            "  待補件：內容過少，尚無法評審",
            "  紅燈：缺少必要章節或計畫不可行",
            "",
            "預期效益總結：",
            "  效率：審核 cycle time 從平均 4 小時/件 → <2 小時/件（>50% 改善）",
            "  規模：Y25 共 356 件案例，預計節省 >793 小時/年",
            "  品質：AI 初篩確保基本規範符合，減少人工重工",
            "  可視化：即時掌握案件狀態，取代 Excel 手工追蹤",
            "  知識累積：每次審核結果自動更新 RAG，系統越用越準",
            "",
            "下一步行動：",
            "  1. 提交本規劃書至 AI 評審中心取得綠燈核准",
            "  2. 5/2 啟動 M1 設計階段",
            "  3. 邀請廠務 SME 盤點歷史審核案例（RAG 知識庫初始資料）",
            "  4. 8/21 全系統上線，正式取代 Excel 作業",
        ],
    ),
]


def _add_slide(prs: Presentation, title: str, body: list[str]) -> None:
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)

    # 標題
    title_shape = slide.shapes.title
    title_shape.text = title
    tf = title_shape.text_frame
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x4D, 0x3C)

    # 內文文字框
    left = prs.slide_width // 12
    top = Emu(1_600_000)
    width = prs.slide_width - 2 * left
    height = prs.slide_height - top - Emu(300_000)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(body):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        run = para.runs[0] if para.runs else para.add_run()
        run.text = line

        # 空行
        if not line.strip():
            run.font.size = Pt(8)
            continue

        # 標題行（不以空格開頭）
        if not line.startswith(" ") and not line.startswith("·") and not line.startswith("S"):
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x4D, 0x3C)
        else:
            run.font.size = Pt(13)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x2B, 0x26, 0x20)


def build() -> None:
    prs = Presentation()
    # 標準寬螢幕 16:9
    prs.slide_width = Emu(9_144_000)
    prs.slide_height = Emu(5_143_500)

    for title, body in SLIDES:
        _add_slide(prs, title, body)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(SLIDES)} slides)")


if __name__ == "__main__":
    build()
