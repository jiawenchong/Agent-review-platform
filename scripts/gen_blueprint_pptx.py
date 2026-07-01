"""產生「AI Agent 開發規劃書」PPT 範本。

給提案人填寫用;上傳後後端會擷取文字 → AI 評審中心判讀綠/紅/待補件。
內容對齊 docs/AGENT_BLUEPRINT_TEMPLATE.md 與 review_document 的必要章節
(目標 / 範圍 / 時程 / 風險 / 資源 / 里程碑)。

用法:python scripts/gen_blueprint_pptx.py
輸出:public/agent-blueprint-template.pptx
     (Vite build 時會複製進 dist/,前端「規劃書評估」頁提供下載;下載檔名顯示為中文)
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

OUT = Path(__file__).resolve().parents[1] / "public" / "agent-blueprint-template.pptx"

# 每張投影片:(標題, [段落…])。標題刻意保留必要章節關鍵字,讓評審能判讀。
SLIDES: list[tuple[str, list[str]]] = [
    (
        "AI Agent 開發規劃書範本(GRAB v1.0)",
        [
            "治理就緒藍圖 — Governance-Ready Agent Blueprint",
            "填寫方式:把每頁方括號 [ ] 的提示換成實際內容,保留章節標題。",
            "Agent 名稱:[例:信用風險評估 Agent]",
            "提案人 / 部門:[姓名] / [部門]",
            "日期 / 版本:[YYYY/MM/DD] / v0.1",
            "評審狀態:(由系統填寫:綠燈 / 紅燈 / 待補件)",
        ],
    ),
    (
        "一、目標(Goal)★必填",
        [
            "一句話定位:[這個 Agent 解決什麼問題、輸出什麼決策]",
            "成功指標(可量化):",
            "  · 治理指標:[例:風險狀態 100% 被偵測,不漏判不誤判]",
            "  · 效率指標:[例:每輪巡查時間較人工降低 >90%]",
            "  · 品質指標:[例:LLM 判斷與人工覆核一致率 ≥90%]",
        ],
    ),
    (
        "二、範圍(Scope)★必填",
        [
            "納入範圍:[本次要做的事項]",
            "排除範圍:[明確不做的,避免範圍蔓延]",
        ],
    ),
    (
        "三、時程與里程碑(Timeline / Milestone)★必填",
        [
            "里程碑 M1 設計 — 預計完成 [YYYY/MM] — 交付物 [...]",
            "里程碑 M2 開發 — 預計完成 [YYYY/MM] — 交付物 [...]",
            "里程碑 M3 驗證 — 預計完成 [YYYY/MM] — 交付物 [...]",
        ],
    ),
    (
        "四、風險與緩解(Risk)★必填",
        [
            "風險:[例:訓練資料時間窗偏移]",
            "影響:[例:模型漂移]",
            "緩解措施:[例:季度 PSI 監控]",
        ],
    ),
    (
        "五、資源需求(Resource)★必填",
        [
            "人力:[角色 × 人數]",
            "運算:[Local GPU / 雲端 / 模型版本]",
            "資料來源:[KANBAN / RAG 知識庫 / 其他]",
        ],
    ),
    (
        "六、系統架構與模組(Architecture)★必填",
        [
            "建議沿用 Closed Loop 四階段,標註沿用 / 新增:",
            "  · Perception 感知/掃描來源:[...](原有 / 新增)",
            "  · Reasoning 判斷邏輯,哪部分用 LLM:[...](原有 / 新增)",
            "  · Action 依判斷分級派發:[...](原有 / 新增)",
            "  · Feedback 驗證迴圈是否閉合:[...](原有 / 新增)",
        ],
    ),
    (
        "七、流程定義(Flow Definition)— 驅動自動流程圖(選填)",
        [
            "每行一個節點,格式:<節點ID> [<類型>] <標籤> -> <下一節點>",
            "類型:start / end / process / decision / io / subprocess",
            "範例:",
            "  S1 [start] 開始",
            "  S2 [process] 接收上傳規劃書 -> S3",
            "  S3 [decision] 文件解析成功? -> 是:S4 | 否:S9",
            "  S4 [process] AI 評審中心判讀 -> S5",
            "  S5 [decision] 評審結果? -> 綠燈:S6 | 紅燈:S9",
            "  S6 [end] 進入 Closed Loop 進度監控",
            "  S9 [end] 退件 · 要求重新提交",
        ],
    ),
    (
        "評審燈號規則(系統如何判讀)",
        [
            "綠燈:涵蓋 ≥4 項必要章節(目標 / 範圍 / 時程 / 風險 / 資源 / 里程碑),內容具體可執行。",
            "待補件:內容過少,無法構成可評審的規劃書。",
            "紅燈:缺少必要章節,或結構不符規劃書要求、計畫不可行。",
            "※ 綠燈才會自動建立專案並進入每 7 天進度監控。",
        ],
    ),
]


def _set_body_text(text_frame, paragraphs: list[str]) -> None:
    text_frame.word_wrap = True
    for i, line in enumerate(paragraphs):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.text = line
        para.font.size = Pt(16)


def build() -> None:
    prs = Presentation()
    title_only = prs.slide_layouts[5]  # Title Only
    for idx, (title, body) in enumerate(SLIDES):
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = title
        # 內文文字框(標題下方)
        left = prs.slide_width // 12
        top = prs.slide_height // 4
        width = prs.slide_width - 2 * left
        height = prs.slide_height - top - (prs.slide_height // 12)
        box = slide.shapes.add_textbox(left, top, width, height)
        _set_body_text(box.text_frame, body)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    build()
